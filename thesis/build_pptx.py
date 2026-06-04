#!/usr/bin/env python3
"""Build the defense presentation (thesis/presentation.pptx).

Mirrors build_pdf.py in spirit: a single reproducible script that turns the
already-written thesis into the deliverable. Uses sample-presentation-2.pptx as
the *style template* (theme, masters, layouts, fonts, 16:9), strips its content
slides, and rebuilds the canonical slide sequence from PRESENTATION-RULES.md.

Verbatim content (мета / задачі / об'єкт / предмет / висновки) is extracted
straight from the thesis markdown so the /grill-presentation word-for-word
cross-check is byte-identical. Diagrams reuse the thesis figures; the control
example uses the fresh presentation-only screenshots in figures/presentation/.
"""
from __future__ import annotations

import re
import struct
import uuid
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
REPO = HERE.parent
TEMPLATE = REPO / "docs/diploma-rules/references/sample-presentation-2.pptx"
FIG = HERE / "figures"
DEMO = FIG / "presentation"
ASSETS = HERE / ".pptx_assets"  # right-sized image copies; figures/ stays immutable
OUT = HERE / "presentation.pptx"

EMU = 914400
# Image-slide titles are set to 28pt (vs the master's 40pt) so even the longest
# "Контрольний приклад: …" stays on ONE line. IMG_TOP=1.9" still clears a 2-line
# title as a safety margin, and IMG_MAX_H keeps the image off the bottom edge.
TITLE_FONT_PT = 28
IMG_TOP = 1.90
IMG_MAX_W = 11.8
IMG_MAX_H = 5.10  # 1.90 + 5.10 = 7.00" of 7.50"
# Keep full source resolution for every image (only cap the longest side to bound
# file size). Quality over size — view in Slideshow / desktop PowerPoint, since
# PowerPoint Online recompresses anything above ~220 PPI on its own side.
IMG_MAX_SIDE = 4000

THEME = "Вебзастосунок для пошуку загублених домашніх тварин"


# --------------------------------------------------------------------------- #
# Verbatim extraction from the thesis markdown                                  #
# --------------------------------------------------------------------------- #
def read(name: str) -> str:
    return (HERE / name).read_text(encoding="utf-8")


def extract_intro() -> dict[str, object]:
    vstup = read("03-vstup.md")
    # Goal = the first sentence of the meta paragraph (the measurable goal itself);
    # the operationalisation/metrics live on the experiment slide. Robust to wording:
    # take the span from "Метою роботи є" up to the tasks connector, first sentence.
    meta_region = vstup[vstup.index("Метою роботи є") : vstup.index("Для досягнення")]
    meta = meta_region.split(". ")[0].strip() + "."
    connector = "Для досягнення поставленої мети необхідно вирішити такі задачі:"
    region = vstup[vstup.index("такі задачі:") + len("такі задачі:") : vstup.index("**Об")]
    tasks = [t.strip().rstrip(";.") for t in re.findall(r"–\s+(.+)", region) if t.strip()]
    # Object / subject are single sentences (no internal periods) — match verbatim.
    obj = re.search(r"Об.єктом роботи є[^.]*\.", vstup).group(0).strip()
    subj = re.search(r"Предметом роботи є[^.]*\.", vstup).group(0).strip()
    return {"meta": meta, "connector": connector, "tasks": tasks, "obj": obj, "subj": subj}


def extract_conclusions() -> list[str]:
    # Committee rule: «не виносити на слайди абзаци тексту роботи». PRESENTATION-RULES §3
    # вимагає, щоб висновки ПЕРЕКАЗУВАЛИ роботу (не дослівні абзаци) і завершувалися
    # формулою «задачі виконано, мету досягнуто». Тож це стислі тези, виведені з
    # 11-zagalni-vysnovky.md, а не вирвані з нього повні речення.
    return [
        "Розроблено вебзастосунок для пошуку загублених тварин через автоматичне "
        "зіставлення оголошень про втрату та про знахідку.",
        "Аналіз аналогів (PawBoost, Pet FBI, спільноти): жоден не зіставляє "
        "автоматично, не ранжує збіги й не приховує контакти до підтвердження.",
        "Ключовий результат — алгоритм зіставлення: гаверсинус для відстані та "
        "лексикографічне ранжування кандидатів; остаточне рішення про збіг — за людиною.",
        "Реалізовано за принципом портів і адаптерів на Node.js, NestJS, React, "
        "Drizzle та PostgreSQL; контакти приховано до підтвердження зв'язку.",
        "Тестування: 130 серверних і 44 клієнтських тести проходять, єдиний конвеєр "
        "перевірки завершується без помилок.",
        "Досягнення мети підтверджено експериментально: істинний збіг у перших трьох "
        "позиціях у всіх запитах (Hit@3 = 1,00), MRR = 0,91; час формування переліку "
        "для 50 000 оголошень ≤ 50 мс, тобто в межах вимоги у 2 с.",
        # Rule §3/§7: mandated closing.
        "Поставлені в роботі задачі виконано повністю, мету роботи досягнуто.",
    ]


# --------------------------------------------------------------------------- #
# Low-level slide helpers                                                       #
# --------------------------------------------------------------------------- #
def png_size(path: Path) -> tuple[int, int]:
    with path.open("rb") as f:
        head = f.read(24)
    # PNG: width/height are big-endian uint32 at byte offset 16.
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def strip_content_slides(prs: Presentation) -> None:
    """Drop the template's content slides — both the sldId entries and the
    underlying slide parts (via their relationships) so they don't linger in
    the saved package as duplicate zip entries."""
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        rId = sld.get(qn("r:id"))
        prs.part.drop_rel(rId)
        lst.remove(sld)


def layout(prs: Presentation, name: str):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(name)


def remove_slide_number(slide) -> None:
    for ph in list(slide.placeholders):
        if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            ph._element.getparent().remove(ph._element)


# Official committee rule: page number lives in the TOP-RIGHT corner. The template
# layout pins its sldNum placeholder bottom-right, so we override the position with
# an explicit xfrm. top=0.15" clears the title band (titles start at ~0.6").
NUM_LEFT_IN = 12.2
NUM_TOP_IN = 0.15
NUM_W_IN = 0.8
NUM_H_IN = 0.4


def _force_top_right(sp) -> None:
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = sp.makeelement(qn("p:spPr"), {})
        sp.append(spPr)
    old = spPr.find(qn("a:xfrm"))
    if old is not None:
        spPr.remove(old)
    xfrm = spPr.makeelement(qn("a:xfrm"), {})
    xfrm.append(xfrm.makeelement(qn("a:off"), {"x": str(int(NUM_LEFT_IN * EMU)), "y": str(int(NUM_TOP_IN * EMU))}))
    xfrm.append(xfrm.makeelement(qn("a:ext"), {"cx": str(int(NUM_W_IN * EMU)), "cy": str(int(NUM_H_IN * EMU))}))
    spPr.insert(0, xfrm)


def add_slide_number(slide, number: int) -> None:
    """Attach an explicit sldNum placeholder (cloned from the layout) carrying a
    slidenum field, so numbering is both rendered and machine-detectable. The
    placeholder is repositioned to the TOP-RIGHT corner per the committee rule."""
    src = None
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            src = ph._element
            break
    if src is None:
        return
    sp = deepcopy(src)
    _force_top_right(sp)
    txBody = sp.find(qn("p:txBody"))
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    p = txBody.makeelement(qn("a:p"), {})
    p.append(p.makeelement(qn("a:pPr"), {"algn": "r"}))
    fld = p.makeelement(qn("a:fld"), {"id": "{%s}" % str(uuid.uuid4()).upper(), "type": "slidenum"})
    t = fld.makeelement(qn("a:t"), {})
    t.text = str(number)
    fld.append(t)
    p.append(fld)
    txBody.append(p)
    slide.shapes._spTree.append(sp)


def _pPr(p):
    return p._p.get_or_add_pPr()


def set_no_bullet(p) -> None:
    pPr = _pPr(p)
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def set_bullet(p) -> None:
    pPr = _pPr(p)
    pPr.set("marL", "342900")
    pPr.set("indent", "-342900")
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))


def fill_body(body, paragraphs: list[dict]) -> None:
    """paragraphs: list of {text, size, bold?, bullet?}."""
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = spec["text"]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(spec["size"])
        run.font.bold = spec.get("bold", False)
        if spec.get("bullet"):
            set_bullet(p)
        else:
            set_no_bullet(p)
        p.space_after = Pt(spec.get("space_after", 6))


def fit_image(image: Path) -> Path:
    """Keep full source resolution; only downscale if the longest side exceeds
    IMG_MAX_SIDE (bounds file size). Never upscale. Diagrams and screenshots both
    stay crisp when zoomed — no shrinking to the small display box."""
    with Image.open(image) as im:
        nw, nh = im.size
        longest = max(nw, nh)
        if longest <= IMG_MAX_SIDE:
            return image
        ASSETS.mkdir(exist_ok=True)
        out = ASSETS / image.name
        scale = IMG_MAX_SIDE / longest
        im.resize((round(nw * scale), round(nh * scale)), Image.LANCZOS).save(
            out, "PNG", optimize=True
        )
        return out


def add_image_slide(prs, title_text: str, image: Path, number: int):
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    slide.shapes.title.text = title_text
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(TITLE_FONT_PT)  # keep long titles on one line
    w_px, h_px = png_size(image)
    ratio = w_px / h_px
    max_w = Emu(int(IMG_MAX_W * EMU))
    max_h = Emu(int(IMG_MAX_H * EMU))
    if max_w / ratio <= max_h:
        width = max_w
        height = Emu(int(max_w / ratio))
    else:
        height = max_h
        width = Emu(int(max_h * ratio))
    src = fit_image(image)
    left = Emu(int((prs.slide_width - width) / 2))
    # Center within the band below the title.
    top = Emu(int(IMG_TOP * EMU + (max_h - height) / 2))
    slide.shapes.add_picture(str(src), left, top, width, height)
    add_slide_number(slide, number)
    return slide


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------- #
# Build                                                                         #
# --------------------------------------------------------------------------- #
def build() -> None:
    intro = extract_intro()
    conclusions = extract_conclusions()

    prs = Presentation(str(TEMPLATE))
    strip_content_slides(prs)

    n = 0  # running slide number (title is unnumbered)

    # 1 — Title -------------------------------------------------------------- #
    s = prs.slides.add_slide(layout(prs, "TITLE"))
    s.shapes.title.text = THEME
    sub = s.placeholders[1].text_frame
    sub.clear()
    sub_lines = [
        ("Кваліфікаційна робота бакалавра", True),
        ("Виконав: ст. гр. АС-222 Ілля Вербанов", False),
        ("Керівник: Надія Беглова, асистент каф. ІПЗ", False),
        ("Національний університет «Одеська політехніка»", False),
        ("Одеса – 2026", False),
    ]
    for i, (line, bold) in enumerate(sub_lines):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.runs[0].font.bold = bold
        p.runs[0].font.size = Pt(22 if bold else 18)
    remove_slide_number(s)
    set_notes(s, NOTES["title"])

    # 2 — Мета і задачі ------------------------------------------------------ #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_AND_BODY"))
    s.shapes.title.text = "Мета і задачі"
    paras = [
        {"text": intro["meta"], "size": 18, "bold": True, "space_after": 4},
        {"text": intro["connector"], "size": 18, "space_after": 2},
    ]
    paras += [{"text": t, "size": 18, "bullet": True, "space_after": 2} for t in intro["tasks"]]
    fill_body(s.placeholders[1], paras)
    # Goal + 9 tasks at the mandated ≥18 pt overrun the layout's 4.87" body box.
    # Use the full content area (just under the title down to ~0.3" of the bottom)
    # so everything fits at 18 pt without dropping any task. Single line spacing.
    # Set ALL FOUR geometry fields: on an inherited placeholder, setting only
    # top/height materializes an xfrm with left=width=0, collapsing the text box
    # to zero width (one char per line). Reuse the layout's inherited left/width.
    body = s.placeholders[1]
    body.left = Emu(int(0.455 * EMU))
    body.width = Emu(int(12.424 * EMU))
    body.top = Emu(int(1.55 * EMU))
    body.height = Emu(int(5.65 * EMU))
    for para in body.text_frame.paragraphs:
        para.line_spacing = 1.0
    add_slide_number(s, n)
    set_notes(s, NOTES["goal"])

    # 3 — Предмет та об'єкт --------------------------------------------------- #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_AND_BODY"))
    s.shapes.title.text = "Об'єкт та предмет"
    fill_body(
        s.placeholders[1],
        [
            {"text": "Об'єкт роботи", "size": 20, "bold": True, "space_after": 4},
            {"text": intro["obj"], "size": 18, "space_after": 12},
            {"text": "Предмет роботи", "size": 20, "bold": True, "space_after": 4},
            {"text": intro["subj"], "size": 18, "space_after": 6},
        ],
    )
    add_slide_number(s, n)
    set_notes(s, NOTES["subject"])

    # 4 — Аналіз аналогів (таблиця) ------------------------------------------ #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    s.shapes.title.text = "Аналіз аналогів"
    header = ["Критерій", "PawBoost", "Pet FBI", "Спільноти", "Власна розробка"]
    # Criteria shortened so each fits ONE line at 18pt in the 4.7" first column
    # (committee rule: шрифт ≥ 18 пт).
    rows = [
        ["Структурована модель оголошення", "+", "+", "–", "+"],
        ["Автоматичне зіставлення оголошень", "–", "–", "–", "+"],
        ["Урахування відстані та дати", "частково", "частково", "–", "+"],
        ["Ранжування кандидатів", "–", "–", "–", "+"],
        ["Підтвердження збігу людиною", "–", "–", "–", "+"],
        ["Контакти приховано до підтвердження", "–", "–", "–", "+"],
        ["Безкоштовність базових функцій", "частково", "+", "+", "+"],
    ]
    n_rows, n_cols = len(rows) + 1, len(header)
    left = Emu(int(0.5 * 914400))
    top = Emu(int(1.5 * 914400))
    width = Emu(int(12.3 * 914400))
    height = Emu(int(5.0 * 914400))
    table = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    table.columns[0].width = Emu(int(4.7 * 914400))
    for c in range(1, n_cols):
        table.columns[c].width = Emu(int(1.9 * 914400))
    for c, txt in enumerate(header):
        cell = table.cell(0, c)
        cell.text = txt
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = txt
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(18)
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    add_slide_number(s, n)
    set_notes(s, NOTES["analogs"])

    # 5–7 — Діаграми --------------------------------------------------------- #
    n += 1
    add_image_slide(prs, "Діаграма варіантів використання системи", DEMO / "diagram-usecase.png", n)
    set_notes(prs.slides[-1], NOTES["usecase"])
    n += 1
    add_image_slide(prs, "Загальна архітектура вебзастосунку", DEMO / "diagram-arch.png", n)
    set_notes(prs.slides[-1], NOTES["arch"])
    n += 1
    add_image_slide(prs, "Схема бази даних («сутність — зв'язок»)", DEMO / "diagram-db.png", n)
    set_notes(prs.slides[-1], NOTES["db"])

    # 8 — Стек технологій ---------------------------------------------------- #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_AND_BODY"))
    s.shapes.title.text = "Стек технологій"
    stack = [
        "Платформа та мова: Node.js 22 LTS, TypeScript 5.6",
        "Серверний фреймворк: NestJS 11 (модульність, впровадження залежностей)",
        "Дані: ORM Drizzle 0.44, PostgreSQL 16",
        "Валідація та обробка помилок: Zod, neverthrow",
        "Безпека: argon2id, JWT, заголовки безпеки, обмеження джерел запитів",
        "Клієнт: React 19, Vite 6, React Router, CSS-модулі",
        "Тестування: Vitest, testcontainers (ефемерна БД у контейнері)",
        "Інфраструктура: монорепозиторій pnpm + Turborepo, Docker",
    ]
    fill_body(
        s.placeholders[1],
        [{"text": t, "size": 18, "bullet": True, "space_after": 4} for t in stack],
    )
    add_slide_number(s, n)
    set_notes(s, NOTES["stack"])

    # 9–13 — Контрольний приклад --------------------------------------------- #
    demo = [
        ("Контрольний приклад: авторизація користувача", "demo-auth.png", "demo_auth"),
        ("Контрольний приклад: публічна стрічка оголошень", "demo-browse.png", "demo_browse"),
        ("Контрольний приклад: форма публікації оголошення", "demo-create.png", "demo_create"),
        ("Контрольний приклад: перелік кандидатів", "demo-candidates.png", "demo_candidates"),
        ("Контрольний приклад: підтверджений зв'язок", "demo-match.png", "demo_match"),
    ]
    for title_text, fname, note_key in demo:
        n += 1
        add_image_slide(prs, title_text, DEMO / fname, n)
        set_notes(prs.slides[-1], NOTES[note_key])

    # 14 — Експериментальне дослідження (таблиця результатів, розділ 6) ------- #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    s.shapes.title.text = "Експериментальне дослідження ефективності зіставлення"
    for run in s.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(24)  # long title on one line
    exp_header = ["Показник", "Значення"]
    exp_rows = [
        ["Hit@1 — істинний збіг на першій позиції", "0,82"],
        ["Hit@3 / Hit@5 — у перших трьох / п'яти", "1,00 / 1,00"],
        ["Середній обернений ранг (MRR)", "0,91"],
        ["Медіанний ранг істинного збігу", "1"],
        ["Час формування переліку (50 000 оголошень)", "≤ 50 мс"],
        ["Скорочення обсягу перегляду vs хронологічний", "≈ 100×"],
    ]
    e_rows, e_cols = len(exp_rows) + 1, len(exp_header)
    e_table = s.shapes.add_table(
        e_rows, e_cols, Emu(int(2.4 * EMU)), Emu(int(1.7 * EMU)),
        Emu(int(8.5 * EMU)), Emu(int(4.4 * EMU)),
    ).table
    e_table.columns[0].width = Emu(int(5.7 * EMU))
    e_table.columns[1].width = Emu(int(2.8 * EMU))
    for c, txt in enumerate(exp_header):
        cell = e_table.cell(0, c)
        cell.text = txt
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(exp_rows, start=1):
        for c, txt in enumerate(row):
            cell = e_table.cell(r, c)
            cell.text = txt
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(18)
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    add_slide_number(s, n)
    set_notes(s, NOTES["experiment"])

    # 15 — Загальні висновки ------------------------------------------------- #
    n += 1
    s = prs.slides.add_slide(layout(prs, "TITLE_AND_BODY"))
    s.shapes.title.text = "Загальні висновки"
    fill_body(
        s.placeholders[1],
        [{"text": b, "size": 18, "bullet": True, "space_after": 4} for b in conclusions],
    )
    # Seven dense bullets at 18 pt overrun the layout's 4.87" body box. Use the full
    # content area, and set ALL FOUR geometry fields (top/height alone would collapse
    # an inherited placeholder to left=width=0 — see slide 2). Single line spacing.
    c_body = s.placeholders[1]
    c_body.left = Emu(int(0.455 * EMU))
    c_body.width = Emu(int(12.424 * EMU))
    c_body.top = Emu(int(1.55 * EMU))
    c_body.height = Emu(int(5.65 * EMU))
    for para in c_body.text_frame.paragraphs:
        para.line_spacing = 1.0
    add_slide_number(s, n)
    set_notes(s, NOTES["conclusions"])

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


# --------------------------------------------------------------------------- #
# Speaker notes (~6 min total). Prose, conversational defense narration.        #
# Slide bodies stay verbatim; only these notes are free prose.                  #
# --------------------------------------------------------------------------- #
NOTES = {
    "title": (
        "Доброго дня. Мене звати Ілля Вербанов, група АС-222. Тема моєї роботи — "
        "вебзастосунок для пошуку загублених домашніх тварин. Керівник роботи — "
        "Надія Беглова. Далі коротко розповім, навіщо цей застосунок і що саме зроблено."
    ),
    "goal": (
        "Коли тварина губиться, вирішують перші дні, а власник і той, хто знайшов "
        "тварину, шукають одне одного вручну в різних чатах і на OLX — оголошення ніде "
        "не зіставляються. Тому мета роботи — підвищити ефективність первинного "
        "зіставлення: розробити вебзастосунок, що сам формує ранжований перелік "
        "ймовірних збігів. Ефективність я міряю двома величинами — якістю ранжування і "
        "часом формування переліку, і досягнення мети підтверджую експериментом. Щоб "
        "дійти до цього, роботу розбито на дев'ять задач: від аналізу області й "
        "аналогів до алгоритму зіставлення, проєктування, реалізації, тестування і "
        "самого експерименту."
    ),
    "subject": (
        "Об'єкт роботи — це процес автоматизації зіставлення оголошень про загублених і "
        "знайдених тварин засобами вебзастосунку, тобто інженерна задача, а не сам "
        "пошук тварин. Предмет — те, що я безпосередньо будую: моделі, алгоритм і "
        "програмні засоби, які структуровано подають оголошення й формують ранжований "
        "перелік ймовірних збігів. Остаточне рішення про збіг приймає людина, а не "
        "система."
    ),
    "analogs": (
        "Я подивився на три типові варіанти: PawBoost, базу Pet FBI і тематичні "
        "спільноти у соцмережах. Усі вони дозволяють опублікувати оголошення, але "
        "жоден не зіставляє втрату зі знахідкою автоматично, не ранжує кандидатів і не "
        "ховає контакти до підтвердження. Саме ці порожні клітинки в останньому "
        "стовпці й стали моєю нішею."
    ),
    "usecase": (
        "Ось основні сценарії користувача. Незалежно від того, власник це чи людина, "
        "яка знайшла тварину, шлях однаковий: опублікувати структуроване оголошення, "
        "переглянути перелік ймовірних збігів і підтвердити або відхилити зв'язок. "
        "Перегляд стрічки доступний і без реєстрації."
    ),
    "arch": (
        "Архітектурно застосунок побудований за принципом портів та адаптерів. "
        "Клієнт на React спілкується із сервером на NestJS через HTTP, а доступ до "
        "PostgreSQL ізольований за портом. Завдяки цьому модуль зв'язків не лізе "
        "напряму в таблицю оголошень, а працює через її порт — межі модулів чіткі."
    ),
    "db": (
        "База навмисно проста — лише три таблиці: користувач, оголошення і зв'язок. "
        "Цього достатньо, щоб описати всю предметну область. Зверніть увагу: єдиний "
        "ідентифікатор користувача є зовнішнім ключем для решти таблиць."
    ),
    "stack": (
        "Щодо інструментів. І клієнт, і сервер — на TypeScript, у одному "
        "монорепозиторії. Сервер на Node.js і NestJS, дані через Drizzle і PostgreSQL. "
        "Помилки повертаю як значення через neverthrow, дані валідую через Zod. Паролі "
        "— argon2id, автентифікація — JWT. Окремо відзначу тести: репозиторні тести "
        "ганяються проти справжньої бази в контейнері через testcontainers, без "
        "імітації драйвера."
    ),
    "demo_auth": (
        "Перейдемо до роботи застосунку. Користувач реєструється або входить — це "
        "звичайна форма входу українською. Після входу стають доступні особисті дії: "
        "подати оголошення, переглянути свої оголошення й збіги."
    ),
    "demo_browse": (
        "Це публічна стрічка оголошень. Її видно й без входу, але контактні дані тут "
        "приховані — їх не показують, доки збіг не підтверджено. Є фільтри за типом і "
        "видом тварини."
    ),
    "demo_create": (
        "Так виглядає публікація оголошення. Форма структурована: тип, вид, кличка, "
        "порода, забарвлення, опис, а далі — місце й дата події. Саме ці поля потім "
        "використовує алгоритм зіставлення."
    ),
    "demo_candidates": (
        "А ось головне — перелік кандидатів для конкретного оголошення про втрату. "
        "Система сама підібрала оголошення про знахідку того ж виду поблизу і "
        "впорядкувала їх: спершу за збігом виду, потім за відстанню, потім за різницею "
        "дат. Власнику лишається переглянути й запропонувати зв'язок."
    ),
    "demo_match": (
        "Коли друга сторона підтверджує зв'язок, контактні дані обох оголошень "
        "взаємно розкриваються — ось ця панель. До цього моменту телефони й пошта були "
        "прихованими. Так ми захищаємо персональні дані й водночас даємо людям зв'язок, "
        "коли збіг справді підтверджено."
    ),
    "experiment": (
        "Щоб показати, що мету досягнуто, я провів експеримент на модельному наборі: "
        "шістдесят пар втрата–знахідка з відомими збігами плюс сто сорок сторонніх "
        "оголошень, усе відтворюється за фіксованим сидом. Істинний збіг потрапляє в "
        "перші три позиції в усіх запитах, на першу — у вісімдесяти двох відсотках, "
        "середній обернений ранг дорівнює нуль цілих дев'яносто одна. Порівняно з "
        "хронологічним переглядом це приблизно стократне скорочення обсягу перегляду. "
        "Час формування переліку для п'ятдесяти тисяч оголошень — до п'ятдесяти "
        "мілісекунд, тобто вимога у дві секунди виконується із запасом."
    ),
    "conclusions": (
        "Підсумую. Я розробив вебзастосунок, який автоматизує найтрудомісткіший етап "
        "пошуку — виявлення ймовірного збігу, і захищає контакти до підтвердження. "
        "Ключовий результат — алгоритм зіставлення на формулі гаверсинуса з "
        "лексикографічним упорядкуванням. Систему перевірено: 130 серверних і 44 "
        "клієнтських тести проходять. А досягнення мети підтверджено експериментально — "
        "істинний збіг у перших трьох позиціях у всіх запитах, час до п'ятдесяти "
        "мілісекунд на п'ятдесяти тисячах оголошень. Поставлені задачі виконано, мету "
        "роботи досягнуто. Дякую за увагу, готовий відповісти на запитання."
    ),
}


if __name__ == "__main__":
    build()
